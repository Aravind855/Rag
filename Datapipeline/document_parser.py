import json
import logging
import re
import threading
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Union

# Optional third-party imports (handled gracefully if missing)
try:
    import pymupdf as fitz
except ImportError:
    try:
        import fitz
    except ImportError:
        fitz = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)


@dataclass
class DocumentElement:
    """Represents a structural unit within a document (e.g. page, slide, sheet, section)."""
    element_type: str  # 'page', 'slide', 'sheet', 'section', 'document'
    index: int  # 1-indexed number
    title: Optional[str] = None
    content: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class ParsedDocument:
    """Dynamic, standardized container for parsed document content and rich metadata."""
    file_path: str
    file_name: str
    file_type: str
    category: str
    content: str
    elements: List[DocumentElement] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)
    error: Optional[str] = None

    @property
    def total_elements(self) -> int:
        return len(self.elements)

    def to_dict(self) -> Dict[str, Any]:
        """Convert ParsedDocument object to dictionary format."""
        return {
            "file_path": self.file_path,
            "file_name": self.file_name,
            "file_type": self.file_type,
            "category": self.category,
            "content": self.content,
            "elements": [
                {
                    "element_type": el.element_type,
                    "index": el.index,
                    "title": el.title,
                    "content": el.content,
                    "metadata": el.metadata,
                }
                for el in self.elements
            ],
            "metadata": self.metadata,
            "error": self.error,
        }


class DocumentParser:
    """Dynamic, extensible document parser.

    Can handle any arbitrary file extension, unknown directory structures,
    and supports runtime registration of custom format handlers.
    """

    def __init__(self):
        # Map file extensions to specialized handler methods
        self._handlers: Dict[str, Callable[[Path, str], ParsedDocument]] = {}
        self._pdf_lock = threading.Lock()
        self._register_default_handlers()

    def register_handler(self, extensions: Union[str, List[str]], handler_fn: Callable[[Path, str], ParsedDocument]):
        """Dynamically registers a custom parser function for given file extension(s).

        Args:
            extensions: Single extension string (e.g. '.xml') or list of extensions (['.html', '.htm']).
            handler_fn: Callable taking (Path, category_str) and returning ParsedDocument.
        """
        ext_list = [extensions] if isinstance(extensions, str) else extensions
        for ext in ext_list:
            clean_ext = ext.lower() if ext.startswith(".") else f".{ext.lower()}"
            self._handlers[clean_ext] = handler_fn
            logger.debug(f"Registered custom parser handler for extension: {clean_ext}")

    def parse(self, file_path: Union[str, Path], category: str = "general") -> ParsedDocument:
        """Parses any document file dynamically regardless of extension or structure.

        Args:
            file_path: Path to the document file.
            category: Category, department, or directory path label.

        Returns:
            ParsedDocument object containing extracted content and structural elements.
        """
        path = Path(file_path).resolve()
        if not path.exists():
            return ParsedDocument(
                file_path=str(path),
                file_name=path.name,
                file_type=path.suffix.lower() or "unknown",
                category=category,
                content="",
                error=f"File not found: {path}",
            )

        file_ext = path.suffix.lower()

        # Check if specialized handler is registered
        handler = self._handlers.get(file_ext)

        if handler:
            try:
                return handler(path, category)
            except Exception as e:
                logger.warning(f"Specialized handler failed for {path.name} ({str(e)}). Retrying with universal text fallback...")

        # Universal Fallback Parser: handles unknown extensions, code, plain text, markup, data files
        return self._parse_universal_fallback(path, category)

    def _register_default_handlers(self):
        """Registers default handlers for standard formats."""
        self.register_handler([".pdf"], self._parse_pdf)
        self.register_handler([".docx", ".doc"], self._parse_docx)
        self.register_handler([".pptx", ".ppt"], self._parse_pptx)
        self.register_handler([".xlsx", ".xls"], self._parse_excel)
        self.register_handler([".csv", ".tsv"], self._parse_csv)
        self.register_handler([".json", ".jsonl"], self._parse_json)
        self.register_handler([".md", ".markdown", ".txt", ".log", ".rst", ".yaml", ".yml", ".ini", ".conf", ".xml", ".html", ".py", ".js"], self._parse_text)

    def _parse_pdf(self, path: Path, category: str) -> ParsedDocument:
        """Parses PDF documents page by page using PyMuPDF (fitz)."""
        if fitz is None:
            logger.info("PyMuPDF unavailable. Falling back to universal text reader for PDF.")
            return self._parse_universal_fallback(path, category)

        with self._pdf_lock:
            elements: List[DocumentElement] = []
            full_text_blocks: List[str] = []

            doc = fitz.open(str(path))
            pdf_metadata = doc.metadata or {}

            for page_num in range(len(doc)):
                page = doc[page_num]
                tabs = page.find_tables()
                blocks = page.get_text("blocks")
                processed_tables = set()
                page_text_parts: List[str] = []
                for b in blocks:
                    block_rect = fitz.Rect(b[:4])
                    block_text = b[4].strip()
                    if not block_text:
                        continue

                    matched_table_idx = None
                    for idx, tab in enumerate(tabs.tables):
                        tab_rect = fitz.Rect(tab.bbox)
                        if block_rect.intersects(tab_rect):
                            matched_table_idx = idx
                            break

                    if matched_table_idx is not None:
                        if matched_table_idx not in processed_tables:
                            tab = tabs.tables[matched_table_idx]
                            try:
                                df = tab.to_pandas()
                                md_table = df.to_markdown(index=False)
                                page_text_parts.append(md_table)
                            except Exception:
                                page_text_parts.append(block_text)
                            processed_tables.add(matched_table_idx)
                    else:
                        page_text_parts.append(block_text)

                for idx, tab in enumerate(tabs.tables):
                    if idx not in processed_tables:
                        try:
                            df = tab.to_pandas()
                            md_table = df.to_markdown(index=False)
                            page_text_parts.append(md_table)
                        except Exception:
                            pass

                text = "\n\n".join(page_text_parts).strip()
                if text:
                    elements.append(
                        DocumentElement(
                            element_type="page",
                            index=page_num + 1,
                            title=f"Page {page_num + 1}",
                            content=text,
                            metadata={
                                "page_number": page_num + 1,
                                "table_count": len(tabs.tables),
                                "has_tables": len(tabs.tables) > 0,
                            },
                        )
                    )
                    full_text_blocks.append(f"--- Page {page_num + 1} ---\n{text}")

            doc.close()

            return ParsedDocument(
                file_path=str(path),
                file_name=path.name,
                file_type=".pdf",
                category=category,
                content="\n\n".join(full_text_blocks),
                elements=elements,
                metadata={
                    "total_pages": len(elements),
                    "author": pdf_metadata.get("author", ""),
                    "title": pdf_metadata.get("title", ""),
                    "creation_date": pdf_metadata.get("creationDate", ""),
                },
            )

    def _parse_docx(self, path: Path, category: str) -> ParsedDocument:
        """Parses DOCX documents including headings, paragraphs, and tables."""
        if docx is None:
            logger.info("python-docx unavailable. Falling back to universal text reader for DOCX.")
            return self._parse_universal_fallback(path, category)

        doc = docx.Document(str(path))
        elements: List[DocumentElement] = []
        full_text_blocks: List[str] = []

        current_heading = "Overview"
        current_section_lines: List[str] = []
        section_idx = 1

        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            if paragraph.style.name.startswith("Heading"):
                if current_section_lines:
                    sec_text = "\n".join(current_section_lines)
                    elements.append(
                        DocumentElement(
                            element_type="section",
                            index=section_idx,
                            title=current_heading,
                            content=sec_text,
                        )
                    )
                    full_text_blocks.append(f"## {current_heading}\n{sec_text}")
                    section_idx += 1
                    current_section_lines = []

                current_heading = text
            else:
                current_section_lines.append(text)

        if current_section_lines:
            sec_text = "\n".join(current_section_lines)
            elements.append(
                DocumentElement(
                    element_type="section",
                    index=section_idx,
                    title=current_heading,
                    content=sec_text,
                )
            )
            full_text_blocks.append(f"## {current_heading}\n{sec_text}")

        # Parse tables
        table_blocks: List[str] = []
        for tbl_idx, table in enumerate(doc.tables, start=1):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                table_data.append(" | ".join(row_data))

            if table_data:
                table_str = "\n".join(table_data)
                table_blocks.append(table_str)
                elements.append(
                    DocumentElement(
                        element_type="table",
                        index=tbl_idx,
                        title=f"Table {tbl_idx}",
                        content=table_str,
                    )
                )

        if table_blocks:
            full_text_blocks.append("--- Tables ---\n" + "\n\n".join(table_blocks))

        return ParsedDocument(
            file_path=str(path),
            file_name=path.name,
            file_type=".docx",
            category=category,
            content="\n\n".join(full_text_blocks),
            elements=elements,
            metadata={
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "section_count": len(elements),
            },
        )

    def _parse_pptx(self, path: Path, category: str) -> ParsedDocument:
        """Parses PPTX presentations slide by slide including shapes, tables, and notes."""
        if Presentation is None:
            logger.info("python-pptx unavailable. Falling back to universal text reader for PPTX.")
            return self._parse_universal_fallback(path, category)

        prs = Presentation(str(path))
        elements: List[DocumentElement] = []
        full_text_blocks: List[str] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_text_parts: List[str] = []
            slide_title = f"Slide {slide_idx}"
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        slide_text_parts.append(" | ".join(row_cells))

            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text_parts.append(f"[Speaker Notes: {notes_text}]")

            slide_content = "\n".join(slide_text_parts)
            if slide_content:
                elements.append(
                    DocumentElement(
                        element_type="slide",
                        index=slide_idx,
                        title=slide_title,
                        content=slide_content,
                        metadata={"has_notes": bool(notes_text)},
                    )
                )
                full_text_blocks.append(f"--- Slide {slide_idx}: {slide_title} ---\n{slide_content}")

        return ParsedDocument(
            file_path=str(path),
            file_name=path.name,
            file_type=".pptx",
            category=category,
            content="\n\n".join(full_text_blocks),
            elements=elements,
            metadata={"total_slides": len(prs.slides)},
        )

    def _format_dataframe(self, df: Any) -> str:
        """Safely converts pandas DataFrame to markdown table or plain text format."""
        try:
            return df.to_markdown(index=False)
        except Exception:
            return df.to_string(index=False)

    def _parse_excel(self, path: Path, category: str) -> ParsedDocument:
        """Parses Excel workbooks (.xlsx, .xls) sheet by sheet."""
        if pd is None or openpyxl is None:
            logger.debug("pandas or openpyxl unavailable. Falling back to universal text reader for Excel.")
            return self._parse_universal_fallback(path, category)

        excel_file = pd.ExcelFile(str(path))
        elements: List[DocumentElement] = []
        full_text_blocks: List[str] = []

        for sheet_idx, sheet_name in enumerate(excel_file.sheet_names, start=1):
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            df = df.dropna(how="all")

            if not df.empty:
                sheet_text = self._format_dataframe(df)
                elements.append(
                    DocumentElement(
                        element_type="sheet",
                        index=sheet_idx,
                        title=sheet_name,
                        content=sheet_text,
                        metadata={"rows": len(df), "columns": len(df.columns)},
                    )
                )
                full_text_blocks.append(f"--- Sheet: {sheet_name} ---\n{sheet_text}")

        return ParsedDocument(
            file_path=str(path),
            file_name=path.name,
            file_type=path.suffix.lower(),
            category=category,
            content="\n\n".join(full_text_blocks),
            elements=elements,
            metadata={
                "total_sheets": len(excel_file.sheet_names),
                "sheet_names": excel_file.sheet_names,
            },
        )

    def _parse_csv(self, path: Path, category: str) -> ParsedDocument:
        """Parses CSV/TSV files into table representations."""
        if pd is None:
            return self._parse_text(path, category)

        sep = "\t" if path.suffix.lower() == ".tsv" else ","
        df = pd.read_csv(str(path), sep=sep)
        df = df.dropna(how="all")
        content = self._format_dataframe(df)

        element = DocumentElement(
            element_type="sheet",
            index=1,
            title=path.stem,
            content=content,
            metadata={"rows": len(df), "columns": len(df.columns)},
        )

        return ParsedDocument(
            file_path=str(path),
            file_name=path.name,
            file_type=path.suffix.lower(),
            category=category,
            content=content,
            elements=[element],
            metadata={"rows": len(df), "columns": len(df.columns)},
        )

    def _parse_json(self, path: Path, category: str) -> ParsedDocument:
        """Parses JSON / JSONL documents dynamically."""
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                data = json.load(f)
            formatted = json.dumps(data, indent=2)
        except Exception:
            return self._parse_text(path, category)

        element = DocumentElement(
            element_type="document",
            index=1,
            title=path.stem,
            content=formatted,
        )

        return ParsedDocument(
            file_path=str(path),
            file_name=path.name,
            file_type=path.suffix.lower(),
            category=category,
            content=formatted,
            elements=[element],
            metadata={"data_type": type(data).__name__},
        )

    def _parse_text(self, path: Path, category: str) -> ParsedDocument:
        """Parses plain text, code, markup, or configuration files."""
        content = self._read_file_text_with_encodings(path)
        lines = content.splitlines()

        element = DocumentElement(
            element_type="document",
            index=1,
            title=path.stem,
            content=content,
            metadata={"line_count": len(lines)},
        )

        return ParsedDocument(
            file_path=str(path),
            file_name=path.name,
            file_type=path.suffix.lower() or ".txt",
            category=category,
            content=content,
            elements=[element],
            metadata={"line_count": len(lines), "char_count": len(content)},
        )

    def _parse_universal_fallback(self, path: Path, category: str) -> ParsedDocument:
        """Universal Fallback Parser for unknown extensions or unhandled formats.

        Tries reading as multi-encoding text. If binary, extracts printable string blocks.
        Guarantees that NO document is dropped or fails silently!
        """
        content = self._read_file_text_with_encodings(path)

        # Check if content looks binary or unreadable
        printable_ratio = len(re.sub(r"[^\x20-\x7E\s]", "", content)) / (len(content) or 1)

        if printable_ratio < 0.6:
            # Binary fallback: extract printable character sequences (strings)
            raw_bytes = path.read_bytes()
            ascii_strings = re.findall(rb"[\x20-\x7E]{4,}", raw_bytes)
            extracted_text = "\n".join(s.decode("ascii", errors="ignore") for s in ascii_strings)
            content = f"[Binary Document Extracted Strings]\n\n{extracted_text}"

        element = DocumentElement(
            element_type="document",
            index=1,
            title=path.stem,
            content=content,
        )

        return ParsedDocument(
            file_path=str(path),
            file_name=path.name,
            file_type=path.suffix.lower() or "unknown",
            category=category,
            content=content,
            elements=[element],
            metadata={"parsed_via": "universal_fallback", "char_count": len(content)},
        )

    def _read_file_text_with_encodings(self, path: Path) -> str:
        """Reads file text trying multiple character encodings."""
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]
        for enc in encodings:
            try:
                with open(path, "r", encoding=enc) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
            except Exception:
                break

        # Fallback reading with character replacement
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


if __name__ == "__main__":
    # Test script for dynamic DocumentParser
    parser = DocumentParser()
    test_dir = Path(r"d:\RAG\Datasource")
    if test_dir.exists():
        print(f"Testing Dynamic DocumentParser on files in {test_dir}...")
        for sample_file in test_dir.rglob("*"):
            if sample_file.is_file() and not sample_file.name.startswith("~$") and sample_file.name != "readme.md":
                res = parser.parse(sample_file, category=sample_file.parent.name)
                print(f"Parsed [{res.category}] {res.file_name} ({res.file_type}): Elements={res.total_elements}, Chars={len(res.content)}, Error={res.error}")
